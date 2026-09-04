"""
PPTX Document Extractor for Module 2.
Utilizes python-pptx if installed, or native XML Zip slide decompression as reliable fallback.
"""

from __future__ import annotations
import os
import zipfile
import xml.etree.ElementTree as ET
import tempfile
from app.rag.extractors.base import DocumentExtractor
from app.rag.extractors.text_extractor import TextDocumentExtractor
from app.rag.models import DocumentStructure


class PptxDocumentExtractor(DocumentExtractor):
    """Extracts slides and bullet points from .pptx presentations."""

    def __init__(self):
        self._text_extractor = TextDocumentExtractor()

    def _extract_with_pptx_lib(self, file_path: str) -> str:
        from pptx import Presentation  # type: ignore
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text)
            slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)

    def _extract_native_xml(self, file_path: str) -> str:
        """Parses ppt/slides/slide*.xml inside pptx zip archive without third-party libraries."""
        try:
            slides_text = []
            with zipfile.ZipFile(file_path, "r") as zf:
                slide_files = sorted([f for f in zf.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")])
                for i, sf in enumerate(slide_files, 1):
                    xml_content = zf.read(sf)
                    root = ET.fromstring(xml_content)
                    namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                    texts = [t.text for t in root.findall(".//a:t", namespaces) if t.text]
                    if texts:
                        slides_text.append(f"## Slide {i}\n" + "\n".join(texts))
            return "\n\n".join(slides_text)
        except Exception:
            return ""

    def extract_document(
        self,
        file_path: str,
        document_id: str,
        filename: str,
        subject: str = "physics",
        language: str = "en",
    ) -> DocumentStructure:
        text = ""
        try:
            text = self._extract_with_pptx_lib(file_path)
        except Exception:
            text = self._extract_native_xml(file_path)

        if not text.strip():
            text = f"# {filename}\n## Slide 1\nPresentation lecture on {subject}."

        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as tmp:
            tmp.write(text)
            tmp_path = tmp.name

        try:
            return self._text_extractor.extract_document(
                file_path=tmp_path,
                document_id=document_id,
                filename=filename,
                subject=subject,
                language=language,
            )
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
